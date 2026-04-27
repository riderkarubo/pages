"""
S8 KPIツリー画像生成スクリプト
出力: KPI_tree_S8.pptx と KPI_tree_S8.png
"""
import sys
sys.path.insert(0, '/Users/issy_s/Library/Python/3.9/lib/python/site-packages')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ========== Fireworkテーマカラー ==========
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT     = RGBColor(0x33, 0x33, 0x33)
C_ACCENT   = RGBColor(0xFA, 0x00, 0x6D)  # マゼンタ #FA006D（MCCM強調色・テキスト/KPI共通）
C_ACCENT2  = RGBColor(0x1B, 0x99, 0x8B)  # 緑
C_GRAY     = RGBColor(0x66, 0x66, 0x66)
C_BORDER   = RGBColor(0xDD, 0xDD, 0xDD)
C_LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)

# ========== スライドサイズ: KPIツリー画像専用 ==========
# 画像配置エリア (Slides側): x=60, y=130, w=600, h=240 → 横600pt × 縦240pt = 比率5:2
# 高解像度で作るため2倍: 1200x480pt → ただしPPTXのEMU単位で扱う
# python-pptxは internal で EMU(914400 = 1inch = 72pt)
SLIDE_W_PT = 600  # pt
SLIDE_H_PT = 240  # pt

prs = Presentation()
prs.slide_width = Pt(SLIDE_W_PT)
prs.slide_height = Pt(SLIDE_H_PT)

# 空レイアウト
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)


def add_rect(left, top, width, height, fill_color, line_color=None, line_width=None, shape=MSO_SHAPE.RECTANGLE):
    """矩形を追加"""
    shp = slide.shapes.add_shape(shape, Pt(left), Pt(top), Pt(width), Pt(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_width is not None:
            shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_text(left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    """テキストボックスを追加"""
    tb = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Noto Sans JP"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_line(x1, y1, x2, y2, color, width=2):
    """直線を追加（細長い矩形で代用）"""
    if abs(x1 - x2) < 0.5:  # 垂直線
        shp = add_rect(x1 - width / 2, min(y1, y2), width, abs(y2 - y1), color)
    else:  # 水平線
        shp = add_rect(min(x1, x2), y1 - width / 2, abs(x2 - x1), width, color)
    shp.line.fill.background()
    return shp


# ========== 描画開始 ==========

# === [1] 上段: 売上ボックス（中央） ===
top_box_w = 200
top_box_h = 50
top_box_x = (SLIDE_W_PT - top_box_w) / 2  # 中央 = 200
top_box_y = 10

add_rect(top_box_x, top_box_y, top_box_w, top_box_h, C_DARK)
add_text(top_box_x, top_box_y, top_box_w, top_box_h, "売上 2,400万円",
         size=18, color=C_WHITE, bold=True)

# === [2] 上段→下段の縦線（=記号まで） ===
center_x = SLIDE_W_PT / 2  # 300
line_top_y = top_box_y + top_box_h  # 60
eq_y = 78  # ＝記号の位置

add_line(center_x, line_top_y, center_x, eq_y, C_BORDER, width=2)

# === [3] ＝記号 ===
add_text(center_x - 15, eq_y, 30, 20, "＝", size=18, color=C_DARK, bold=True)

# === [4] 下段3カード ===
card_y = 110
card_h = 120
card_w = 175
card_gap = 12

# 3カード合計幅 = 3*175 + 2*12 = 549、中央配置: x_start = (600-549)/2 = 25.5
total_cards_w = 3 * card_w + 2 * card_gap
card_x_start = (SLIDE_W_PT - total_cards_w) / 2

# カード上の分岐線（横方向）
# 左カード中央: card_x_start + card_w/2
# 右カード中央: card_x_start + 2*(card_w + card_gap) + card_w/2
left_center_x  = card_x_start + card_w / 2
mid_center_x   = card_x_start + (card_w + card_gap) + card_w / 2
right_center_x = card_x_start + 2 * (card_w + card_gap) + card_w / 2

branch_y = 96  # 分岐線の高さ
add_line(left_center_x, branch_y, right_center_x, branch_y, C_BORDER, width=2)
# 各カードへの縦線
add_line(left_center_x,  branch_y, left_center_x,  card_y, C_BORDER, width=2)
add_line(mid_center_x,   branch_y, mid_center_x,   card_y, C_BORDER, width=2)
add_line(right_center_x, branch_y, right_center_x, card_y, C_BORDER, width=2)
# ＝記号からのつなぎ縦線
add_line(center_x, eq_y + 18, center_x, branch_y, C_BORDER, width=2)

cards = [
    {
        "x": card_x_start,
        "role": "貴社",
        "role_color": C_GRAY,
        "label": "商談件数",
        "value": "貴社で設計",
        "value_color": C_DARK,
        "note": "営業チャネル × アプローチ",
        "top_color": C_GRAY,
    },
    {
        "x": card_x_start + (card_w + card_gap),
        "role": "Firework",
        "role_color": C_ACCENT,
        "label": "受注率",
        "value": "↑ 押し上げ",
        "value_color": C_ACCENT,
        "note": "コンテンツ・実績の説得力",
        "top_color": C_ACCENT,
    },
    {
        "x": card_x_start + 2 * (card_w + card_gap),
        "role": "Firework",
        "role_color": C_ACCENT2,
        "label": "受注単価",
        "value": "↑ 押し上げ",
        "value_color": C_ACCENT2,
        "note": "メディアパワー・差別化",
        "top_color": C_ACCENT2,
    },
]

for i, c in enumerate(cards):
    # カード本体（白背景・グレー枠線）
    shp = add_rect(c["x"], card_y, card_w, card_h, C_WHITE, line_color=C_BORDER, line_width=1)
    # 上ボーダー（4pt）
    add_rect(c["x"], card_y, card_w, 4, c["top_color"])
    # 役割タグ（左上小さく）
    add_text(c["x"] + 10, card_y + 8, card_w - 20, 12, c["role"],
             size=8, color=c["role_color"], bold=True, align=PP_ALIGN.LEFT)
    # ラベル
    add_text(c["x"] + 10, card_y + 26, card_w - 20, 22, c["label"],
             size=14, color=C_DARK, bold=True)
    # 値（強調）
    add_text(c["x"] + 10, card_y + 56, card_w - 20, 30, c["value"],
             size=18, color=c["value_color"], bold=True)
    # 注釈
    add_text(c["x"] + 10, card_y + 92, card_w - 20, 22, c["note"],
             size=10, color=C_GRAY)

    # ×記号（カード間）
    if i < 2:
        x_mark_x = c["x"] + card_w
        add_text(x_mark_x, card_y + card_h / 2 - 12, card_gap, 24, "×",
                 size=16, color=C_DARK, bold=True)

# ========== 保存 ==========
out_dir = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.join(out_dir, "KPI_tree_S8.pptx")
prs.save(pptx_path)
print(f"✅ PPTX 保存: {pptx_path}")
print(f"   スライドサイズ: {SLIDE_W_PT}pt × {SLIDE_H_PT}pt（5:2）")
